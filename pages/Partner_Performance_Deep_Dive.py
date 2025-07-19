import streamlit as st
import pandas as pd
import numpy as np
import plotly.graph_objects as go
import plotly.express as px
import plotly.figure_factory as ff
from sklearn.ensemble import IsolationForest
from prophet import Prophet
from pptx import Presentation
from pptx.util import Inches
import io
import warnings

# --- SUPPRESS DEPRECATION WARNINGS FOR CLEANER PRESENTATION ---
warnings.filterwarnings("ignore", category=DeprecationWarning)
warnings.filterwarnings("ignore", category=FutureWarning)

# --- ROBUST STATE CHECK ---
if 'app_data' not in st.session_state:
    st.error("Application data not loaded. Please return to the 'Global Command Center' home page to initialize the app.")
    st.stop()

# Unpack data
app_data = st.session_state['app_data']
partners = app_data['partners']
batches = app_data['batches']
deviations = app_data['deviations']
micro_data = app_data['micro_data']

# --- UI RENDER ---
st.markdown("# 📈 Partner Performance Deep Dive")
st.markdown("This module provides a detailed analytical view of a single CMO/CTO's performance, ideal for troubleshooting issues or preparing for a Quarterly Business Review (QBR).")

selected_partner = st.selectbox("Select a Partner to Analyze", partners['Partner'].unique())
partner_info = partners[partners['Partner'] == selected_partner].iloc[0]
partner_batches = batches[batches['Partner'] == selected_partner]
partner_devs = deviations[deviations['Partner'] == selected_partner]
partner_micro = micro_data[micro_data['Partner'] == selected_partner]

st.header(f"Performance Analysis for: {selected_partner}")
st.caption(f"**Partner Type:** {partner_info['Type']} | **Specialty:** {partner_info['Specialty']}")


tab_perf, tab_micro, tab_predict, tab_report = st.tabs(["Performance Analytics (SPC/Cpk)", "🔬 Microbiology", "🔮 Predictive Analytics", "📋 QBR Report Generation"])

with tab_perf:
    st.subheader("Key Performance & Quality Metrics")
    col1, col2, col_empty, col3 = st.columns([1,1,0.5,1])
    with col1:
        on_time_rate = (partner_batches['Actual_TAT'] <= partner_batches['TAT_SLA']).mean() * 100 if not partner_batches.empty else 100
        st.metric("On-Time Release Rate", f"{on_time_rate:.1f}%")
    with col2:
        oos_rate = (partner_devs['Type'] == 'OOS').mean() * 100 if not partner_devs.empty else 0
        st.metric("Out-of-Specification (OOS) Rate", f"{oos_rate:.1f}%")
    with col3:
        open_capas = partner_devs[(partner_devs['Status'] == 'CAPA Plan')].shape[0]
        st.metric("Open CAPAs", open_capas)
    
    st.divider()
    
    col1, col2 = st.columns(2)
    with col1:
        st.subheader("Process Capability (Cpk)")
        st.markdown("- **What:** A histogram of a **Critical Quality Attribute (CQA)** measured by this partner, plotted against engineering specifications. \n- **Why (Actionability):** This core **Six Sigma** metric proves if a partner's process is *capable* of reliably meeting our requirements. A **Cpk < 1.33** is a data-driven trigger to demand process improvement per **AS9145 (APQP)** principles, preventing future OOS events.")
        
        usl, lsl = 99.8, 98.0
        mu, sigma = np.random.uniform(98.5, 99.5), np.random.uniform(0.1, 0.4)
        process_data = np.random.normal(mu, sigma, 200)
        
        cpu = (usl - mu) / (3 * sigma)
        cpl = (mu - lsl) / (3 * sigma)
        cpk = min(cpu, cpl)

        fig_cpk = ff.create_distplot([process_data], ['Purity Data'], show_hist=True, show_rug=False, colors=['#444e86'])
        fig_cpk.add_vline(x=usl, line=dict(dash="dash", color="red"), name="USL", annotation_text="USL")
        fig_cpk.add_vline(x=lsl, line=dict(dash="dash", color="red"), name="LSL", annotation_text="LSL")
        fig_cpk.update_layout(title=f"Process Capability: Purity by HPLC (Cpk = {cpk:.2f})")
        st.plotly_chart(fig_cpk, use_container_width=True)

        if cpk < 1.33:
            st.error(f"**Action Required:** Cpk of {cpk:.2f} is below the target of 1.33. This indicates a high risk of producing non-conforming parts. A process improvement plan is required.", icon="🚨")
        else:
            st.success(f"**Process Capable:** Cpk of {cpk:.2f} meets the target of 1.33.", icon="✅")

    with col2:
        st.subheader("Turnaround Time (TAT) Performance")
        st.markdown("- **Why (Actionability):** This provides objective, visual evidence to 'hold partners accountable'. A distribution heavily skewed to the right of the SLA line is clear justification for escalating performance issues during a QBR.")
        if not partner_batches.empty:
            sla = partner_batches['TAT_SLA'].iloc[0]
            fig_hist = px.histogram(partner_batches, x="Actual_TAT", nbins=15, title=f"Testing TAT Distribution (SLA: {sla} days)", color_discrete_sequence=['#955196'])
            fig_hist.add_vline(x=sla, line_dash="dash", line_color="red", annotation_text="SLA")
            st.plotly_chart(fig_hist, use_container_width=True)
        else:
            st.info("No batch data available for this partner to analyze TAT.")

with tab_micro:
    st.header("Microbiology & Environmental Monitoring Trends")
    st.markdown("This section provides specific oversight for microbiological testing and environmental monitoring (EM), critical for ensuring the safety and sterility of our injectable drug products.")
    
    # --- SME ENHANCEMENT: Context-Aware Display ---
    # Check if the selected partner actually performs microbiology testing.
    if not partner_micro.empty:
        st.subheader("Environmental Monitoring (EM) Control Chart")
        st.markdown("- **What:** A c-Chart for microbial counts from Grade B cleanroom surfaces. It plots the number of Colony Forming Units (CFU) per plate against statistical Alert and Action limits. \n- **Why:** This is a primary tool for monitoring the state of control of an aseptic manufacturing environment. A point exceeding the Action Limit is a serious event that requires immediate investigation and potential production holds. \n- **Regulatory:** Aligns with **FDA aseptic processing guidelines** and **EU GMP Annex 1**.")
        em_data = partner_micro[partner_micro['Test'] == 'EM Grade B (CFU/plate)']
        
        fig_em = go.Figure()
        fig_em.add_trace(go.Scatter(x=em_data['Date'], y=em_data['Result'], mode='lines+markers', name='CFU Count', line=dict(color='royalblue')))
        fig_em.add_hline(y=10, line=dict(color="orange", dash="dash"), name="Alert Limit")
        fig_em.add_hline(y=20, line=dict(color="red", dash="dash"), name="Action Limit")
        
        excursions = em_data[em_data['Result'] > 10]
        fig_em.add_trace(go.Scatter(x=excursions['Date'], y=excursions['Result'], mode='markers', name='Excursions', marker=dict(color='red', size=12, symbol='x')))

        fig_em.update_layout(title="EM Trend: Grade B Surface Monitoring", yaxis_title="CFU/plate", xaxis_title="Date", legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1))
        st.plotly_chart(fig_em, use_container_width=True)
    else:
        # Provide an intelligent, context-specific message.
        st.info(f"**Not Applicable:** Microbiology and Environmental Monitoring data is not generated by this supplier type (`{partner_info['Specialty']}`). Please select a CTO with a specialty in 'Analytical & Micro' or a CMO with a specialty in 'Drug Product Fill/Finish' to view this data.", icon="ℹ️")

with tab_predict:
    st.header("Predictive Analytics")
    st.subheader("Future Performance Forecast (Prophet)")
    st.markdown("- **Why:** This shifts the SQE role from reactive to proactive. The forecast warns of potential TAT degradation, allowing for preventive actions.")
    
    @st.cache_data
    def run_prophet_forecast(data, periods=12): # Forecasting 12 weeks
        if data.empty or len(data) < 2: return None, None
        prophet_df = data[['Date_Created', 'Actual_TAT']].rename(columns={'Date_Created': 'ds', 'Actual_TAT': 'y'})
        prophet_df = prophet_df.groupby('ds').mean().reset_index()
        m = Prophet(daily_seasonality=False, weekly_seasonality=True, yearly_seasonality=False)
        m.fit(prophet_df)
        future = m.make_future_dataframe(periods=periods, freq='W')
        return m, m.predict(future)

    model_prophet, forecast = run_prophet_forecast(partner_batches)
    
    if forecast is not None:
        fig_forecast = go.Figure()
        fig_forecast.add_trace(go.Scatter(x=forecast['ds'], y=forecast['yhat'], mode='lines', name='Forecast TAT', line=dict(color='navy', dash='dash')))
        st.plotly_chart(fig_forecast, use_container_width=True)
    else:
        st.info(f"Not enough data to generate a forecast for {selected_partner}.")

with tab_report:
    st.header("QBR Report Generation")
    st.markdown("- **Why (Actionability):** This feature directly addresses the **'excellent written and verbal communication'** requirement by automating the creation of a standardized, data-driven Quarterly Business Review (QBR) deck.")
    
    if st.button("Generate QBR PowerPoint Report"):
        with st.spinner("Creating QBR Deck..."):
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0]); slide.shapes.title.text = f"Quarterly Business Review: {selected_partner}"; slide.placeholders[1].text = f"Avidity Biosciences QC\nReport Generated: {pd.Timestamp.now().strftime('%Y-%m-%d')}"
            
            slide = prs.slides.add_slide(prs.slide_layouts[5]); slide.shapes.title.text = "Key Performance Metrics"
            
            try:
                if 'fig_cpk' in locals():
                    cpk_img_bytes = fig_cpk.to_image(format="png", width=800, height=400)
                    slide.shapes.add_picture(io.BytesIO(cpk_img_bytes), Inches(0.5), Inches(1.5), width=Inches(4.5))
            except Exception as e: st.warning(f"Could not generate Cpk chart image: {e}")

            try:
                if 'fig_hist' in locals() and not partner_batches.empty:
                    hist_img_bytes = fig_hist.to_image(format="png", width=800, height=400)
                    slide.shapes.add_picture(io.BytesIO(hist_img_bytes), Inches(5.0), Inches(1.5), width=Inches(4.5))
            except Exception as e: st.warning(f"Could not generate TAT chart image: {e}")

            ppt_buffer = io.BytesIO()
            prs.save(ppt_buffer)
            ppt_buffer.seek(0)
            
            st.success("Report generated successfully!")
            st.download_button(
                label="Download QBR Deck (.pptx)",
                data=ppt_buffer,
                file_name=f"QBR_{selected_partner.replace(' ', '_')}_{pd.Timestamp.now().strftime('%Y%m%d')}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
